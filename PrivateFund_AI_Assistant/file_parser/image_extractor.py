"""从上传材料中提取可直接放入客户报告的原始业绩图片。"""

from __future__ import annotations

import hashlib
import shutil
import tempfile
import zipfile
from io import BytesIO
from pathlib import Path, PurePosixPath

import fitz
from docx import Document
from openpyxl import load_workbook
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from config.config import SUPPORTED_EXTENSIONS
from models import SourceImage

from .zip_parser import decode_zip_member_name


PERFORMANCE_HINTS = (
    "净值",
    "超额",
    "业绩",
    "收益",
    "回撤",
    "夏普",
    "表现",
    "指数增强",
    "累计",
)
MAX_IMAGES_PER_FILE = 24
MAX_PDF_PAGES = 8
MAX_IMAGE_BYTES = 8 * 1024 * 1024
MIN_WIDTH = 320
MIN_HEIGHT = 150
MIN_PIXELS = 90_000


def extract_source_images(file_paths: list[Path]) -> tuple[list[SourceImage], list[str]]:
    """批量提取图片；单个文件失败不影响文字解析和其他材料。"""
    images: list[SourceImage] = []
    errors: list[str] = []
    for file_path in file_paths:
        try:
            images.extend(_extract_file(file_path))
        except Exception as exc:
            errors.append(f"{file_path.name}（图片提取）：{exc}")
    return assign_image_ids(images), errors


def assign_image_ids(images: list[SourceImage]) -> list[SourceImage]:
    """按来源与图片内容稳定排序并分配内部图片编号。"""
    unique: dict[tuple[str, str, str], SourceImage] = {}
    for item in images:
        digest = hashlib.sha256(item.image_data).hexdigest()
        key = (item.source_file.casefold(), item.source_page or "", digest)
        unique.setdefault(key, item)
    ordered = sorted(
        unique.values(),
        key=lambda item: (
            item.source_file.casefold(),
            (item.source_page or "").casefold(),
            hashlib.sha256(item.image_data).hexdigest(),
        ),
    )
    for index, item in enumerate(ordered, start=1):
        item.image_id = f"I{index:04d}"
    return ordered


def _extract_file(file_path: Path) -> list[SourceImage]:
    suffix = file_path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        return []
    if suffix == ".zip":
        return _extract_zip(file_path)
    if suffix == ".docx":
        return _extract_docx(file_path)
    if suffix == ".pptx":
        return _extract_pptx(file_path)
    if suffix in {".xlsx", ".xlsm"}:
        return _extract_excel(file_path)
    if suffix == ".pdf":
        return _extract_pdf(file_path)
    return []


def _extract_docx(file_path: Path) -> list[SourceImage]:
    document = Document(str(file_path))
    paragraphs = document.paragraphs
    results: list[SourceImage] = []
    used_relationships: set[str] = set()
    for index, paragraph in enumerate(paragraphs):
        relationship_ids = paragraph._p.xpath(".//a:blip/@r:embed")
        if not relationship_ids:
            continue
        context = _nearby_paragraph_text(paragraphs, index)
        for relationship_id in relationship_ids:
            used_relationships.add(relationship_id)
            part = document.part.related_parts.get(relationship_id)
            if part is None:
                continue
            item = _source_image(
                part.blob,
                file_path.name,
                f"段落 {index + 1}",
                context,
                "Word",
            )
            if item:
                results.append(item)
    full_context = "\n".join(paragraph.text.strip() for paragraph in paragraphs if paragraph.text.strip())[:12_000]
    for relationship_id, part in document.part.related_parts.items():
        if relationship_id in used_relationships or not getattr(part, "content_type", "").startswith("image/"):
            continue
        item = _source_image(part.blob, file_path.name, "嵌入图片", full_context, "Word")
        if item:
            results.append(item)
    return results[:MAX_IMAGES_PER_FILE]


def _extract_pptx(file_path: Path) -> list[SourceImage]:
    presentation = Presentation(str(file_path))
    results: list[SourceImage] = []
    for page_number, slide in enumerate(presentation.slides, start=1):
        text = "\n".join(
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text and shape.text.strip()
        )
        for shape in _walk_shapes(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob = shape.image.blob
            except (AttributeError, ValueError):
                continue
            item = _source_image(blob, file_path.name, f"第 {page_number} 页", text, "PPT")
            if item:
                results.append(item)
        if len(results) >= MAX_IMAGES_PER_FILE:
            break
    return results[:MAX_IMAGES_PER_FILE]


def _walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)


def _extract_pdf(file_path: Path) -> list[SourceImage]:
    results: list[SourceImage] = []
    with fitz.open(str(file_path)) as document:
        candidates: list[tuple[int, int, str]] = []
        for page_index, page in enumerate(document):
            context = page.get_text("text").strip()
            score = _performance_text_score(context)
            if score >= 2:
                candidates.append((score, page_index, context))
        candidates.sort(key=lambda item: (-item[0], item[1]))
        for _, page_index, context in candidates[:MAX_PDF_PAGES]:
            page = document[page_index]
            pixmap = page.get_pixmap(matrix=fitz.Matrix(1.35, 1.35), alpha=False)
            item = _source_image(
                pixmap.tobytes("png"),
                file_path.name,
                f"第 {page_index + 1} 页（原页截图）",
                context,
                "PDF",
                allow_page=True,
            )
            if item:
                results.append(item)
    return results


def _extract_excel(file_path: Path) -> list[SourceImage]:
    workbook = load_workbook(str(file_path), read_only=False, data_only=True)
    results: list[SourceImage] = []
    try:
        for worksheet in workbook.worksheets:
            values: list[str] = []
            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.value is not None and str(cell.value).strip():
                        values.append(str(cell.value).strip())
                    if len(values) >= 240:
                        break
                if len(values) >= 240:
                    break
            context = " | ".join(values)
            for image in getattr(worksheet, "_images", []):
                try:
                    blob = image._data()
                except (AttributeError, OSError, ValueError):
                    continue
                item = _source_image(
                    blob,
                    file_path.name,
                    f"工作表“{worksheet.title}”",
                    context,
                    "Excel",
                )
                if item:
                    results.append(item)
            if len(results) >= MAX_IMAGES_PER_FILE:
                break
    finally:
        workbook.close()
    return results[:MAX_IMAGES_PER_FILE]


def _extract_zip(file_path: Path) -> list[SourceImage]:
    results: list[SourceImage] = []
    with tempfile.TemporaryDirectory(prefix="private_fund_images_") as temp_dir:
        target = Path(temp_dir)
        with zipfile.ZipFile(file_path) as archive:
            for member in archive.infolist():
                decoded_name = decode_zip_member_name(member.filename)
                relative_path = PurePosixPath(decoded_name.replace("\\", "/"))
                if relative_path.is_absolute() or ".." in relative_path.parts:
                    raise ValueError(f"压缩包包含不安全路径：{decoded_name}")
                member_path = (target / Path(*relative_path.parts)).resolve()
                if target.resolve() not in member_path.parents and member_path != target.resolve():
                    raise ValueError(f"压缩包包含不安全路径：{decoded_name}")
                if member.is_dir():
                    member_path.mkdir(parents=True, exist_ok=True)
                    continue
                member_path.parent.mkdir(parents=True, exist_ok=True)
                with archive.open(member) as source, member_path.open("wb") as destination:
                    shutil.copyfileobj(source, destination)
        for nested_file in sorted(target.rglob("*")):
            if not nested_file.is_file():
                continue
            try:
                nested_images = _extract_file(nested_file)
            except (ValueError, OSError, zipfile.BadZipFile):
                continue
            relative_name = nested_file.relative_to(target).as_posix()
            for item in nested_images:
                item.source_file = f"{file_path.name} / {relative_name}"
            results.extend(nested_images)
            if len(results) >= MAX_IMAGES_PER_FILE:
                break
    return results[:MAX_IMAGES_PER_FILE]


def _nearby_paragraph_text(paragraphs, index: int) -> str:
    start = max(0, index - 3)
    end = min(len(paragraphs), index + 4)
    return "\n".join(paragraphs[position].text.strip() for position in range(start, end) if paragraphs[position].text.strip())


def _performance_text_score(text: str) -> int:
    lowered = text.casefold()
    return sum(1 for hint in PERFORMANCE_HINTS if hint.casefold() in lowered)


def _source_image(
    data: bytes,
    source_file: str,
    source_page: str,
    context_text: str,
    file_type: str,
    *,
    allow_page: bool = False,
) -> SourceImage | None:
    if not data or len(data) > MAX_IMAGE_BYTES:
        return None
    try:
        with Image.open(BytesIO(data)) as opened:
            opened.load()
            width, height = opened.size
            if not allow_page and (width < MIN_WIDTH or height < MIN_HEIGHT or width * height < MIN_PIXELS):
                return None
            image = opened.convert("RGBA") if opened.mode in {"RGBA", "LA", "P"} else opened.convert("RGB")
            image.thumbnail((2200, 1800), Image.Resampling.LANCZOS)
            width, height = image.size
            buffer = BytesIO()
            image.save(buffer, format="PNG", optimize=True)
    except (UnidentifiedImageError, OSError, ValueError):
        return None
    return SourceImage(
        image_data=buffer.getvalue(),
        source_file=source_file,
        source_page=source_page,
        context_text=context_text[:16_000],
        width=width,
        height=height,
        file_type=file_type,
    )
